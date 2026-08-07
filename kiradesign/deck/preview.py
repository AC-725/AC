#!/usr/bin/env python3
"""
Render a .pptx to images by reading its real geometry with python-pptx.

LibreOffice will not run in this sandbox (it fails on a one-textbox deck),
so this exists purely so the slides can actually be looked at before they go
in front of a client. Fonts are metric-compatible substitutes, so text extents
are close enough to catch overflow, overlap and out-of-bounds placement.
"""
import sys, os, io
from pptx import Presentation
from pptx.util import Emu
from PIL import Image, ImageDraw, ImageFont

EMU_IN = 914400.0
DPI = 110

F = "/usr/share/fonts/truetype/"
FONTS = {
    ("sans", False, False): F + "liberation/LiberationSans-Regular.ttf",
    ("sans", True,  False): F + "liberation/LiberationSans-Bold.ttf",
    ("sans", False, True):  F + "liberation/LiberationSans-Italic.ttf",
    ("sans", True,  True):  F + "liberation/LiberationSans-BoldItalic.ttf",
    ("serif", False, False): F + "liberation/LiberationSerif-Regular.ttf",
    ("serif", True,  False): F + "liberation/LiberationSerif-Bold.ttf",
    ("serif", False, True):  F + "liberation/LiberationSerif-Italic.ttf",
    ("serif", True,  True):  F + "liberation/LiberationSerif-BoldItalic.ttf",
}
CJK = F + "wqy/wqy-zenhei.ttc"
if not os.path.exists(CJK):
    CJK = F + "fonts-japanese-gothic.ttf"

SERIF_NAMES = {"cambria", "times new roman", "century schoolbook",
               "bookman old style", "georgia", "garamond"}


def px(emu):
    return int(round(emu / EMU_IN * DPI))


def font_for(name, bold, italic, size_pt):
    fam = "serif" if (name or "").lower() in SERIF_NAMES else "sans"
    path = FONTS[(fam, bool(bold), bool(italic))]
    return ImageFont.truetype(path, max(6, int(round(size_pt * DPI / 72.0))))


def cjk_font(size_pt):
    try:
        return ImageFont.truetype(CJK, max(6, int(round(size_pt * DPI / 72.0))))
    except Exception:
        return None


def has_cjk(s):
    return any('　' <= c <= '鿿' or '＀' <= c <= '￯' for c in s)


def shape_fill(shape):
    """(rgb, alpha 0..1) or None. python-pptx exposes no alpha, so read the XML."""
    try:
        f = shape.fill
        if f.type is None or f.type != 1:
            return None
        rgb = f.fore_color.rgb
        col = (rgb[0], rgb[1], rgb[2]) if not isinstance(rgb, str) else tuple(
            int(str(rgb)[i:i + 2], 16) for i in (0, 2, 4))
        alpha = 1.0
        el = shape.fill._xPr.find(
            '{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill')
        if el is not None:
            clr = el.find('{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr')
            if clr is not None:
                a = clr.find('{http://schemas.openxmlformats.org/drawingml/2006/main}alpha')
                if a is not None:
                    alpha = int(a.get('val')) / 100000.0
        return col, alpha
    except Exception:
        return None


def run_color(run, default=(0, 0, 0)):
    try:
        c = run.font.color
        if c and c.type is not None and c.rgb is not None:
            v = str(c.rgb)
            return tuple(int(v[i:i + 2], 16) for i in (0, 2, 4))
    except Exception:
        pass
    return default


def wrap(draw, text, fnt, max_w):
    out = []
    for hard in text.split("\n"):
        line = ""
        for word in hard.split(" "):
            trial = (line + " " + word).strip()
            if draw.textlength(trial, font=fnt) <= max_w or not line:
                line = trial
            else:
                out.append(line)
                line = word
        # keep edge whitespace — it is exactly what we are QA-ing for
        if hard.endswith(" ") and line:
            line += " "
        out.append(line)
    return out


def render(path, outdir):
    prs = Presentation(path)
    W, H = px(prs.slide_width), px(prs.slide_height)
    os.makedirs(outdir, exist_ok=True)
    written, warnings = [], []

    for idx, slide in enumerate(prs.slides, 1):
        img = Image.new("RGB", (W, H), (255, 255, 255))

        try:
            bg = slide.background.fill
            if bg.type == 1:
                v = str(bg.fore_color.rgb)
                img.paste(tuple(int(v[i:i + 2], 16) for i in (0, 2, 4)), (0, 0, W, H))
        except Exception:
            pass

        d = ImageDraw.Draw(img, "RGBA")

        for shape in slide.shapes:
            if shape.left is None:
                continue
            x, y = px(shape.left), px(shape.top)
            w, h = px(shape.width), px(shape.height)

            if shape.shape_type == 13:                       # PICTURE
                try:
                    pic = Image.open(io.BytesIO(shape.image.blob)).convert("RGB")
                    sw, sh = pic.size
                    sc = max(w / sw, h / sh)                 # emulate cover
                    pic = pic.resize((max(1, int(sw * sc)), max(1, int(sh * sc))), Image.LANCZOS)
                    ox, oy = (pic.width - w) // 2, (pic.height - h) // 2
                    img.paste(pic.crop((ox, oy, ox + w, oy + h)), (x, y))
                    d = ImageDraw.Draw(img, "RGBA")
                except Exception as e:
                    warnings.append(f"slide {idx}: picture failed — {e}")
                continue

            fill = shape_fill(shape)
            if fill:
                (r, g, b), a = fill
                d.rectangle([x, y, x + w, y + h], fill=(r, g, b, int(255 * a)))

            if not shape.has_text_frame:
                continue

            cy = y + px(45720)                               # default 0.05" inset
            for para in shape.text_frame.paragraphs:
                runs = [r for r in para.runs if r.text]
                if not runs:
                    cy += 6
                    continue
                sizes = [r.font.size.pt for r in runs if r.font.size]
                base = max(sizes) if sizes else 18
                line_h = int(base * DPI / 72.0 * 1.22)

                cx = x + px(91440)
                align = str(para.alignment or "")
                if "RIGHT" in align:
                    total = sum(d.textlength(r.text, font=font_for(
                        r.font.name, r.font.bold, r.font.italic,
                        r.font.size.pt if r.font.size else base)) for r in runs)
                    cx = x + w - px(91440) - total

                for r in runs:
                    sz = r.font.size.pt if r.font.size else base
                    fnt = font_for(r.font.name, r.font.bold, r.font.italic, sz)
                    col = run_color(r, (20, 20, 20))
                    # first line wraps into whatever is left on the current
                    # line; continuation lines get the full box width
                    first_avail = max(40, x + w - px(91440) - cx)
                    full_avail  = max(40, w - px(182880))
                    lines = wrap(d, r.text, fnt, first_avail)
                    if len(lines) > 1:            # re-wrap the tail at full width
                        head, rest = lines[0], " ".join(lines[1:])
                        lines = [head] + wrap(d, rest, fnt, full_avail)
                    for i, line in enumerate(lines):
                        if i:
                            cy += line_h
                            cx = x + px(91440)
                        use = cjk_font(sz) if has_cjk(line) and cjk_font(sz) else fnt
                        d.text((cx, cy), line, font=use, fill=col)
                        cx += d.textlength(line, font=use)
                cy += line_h

                if cy > y + h + 4:
                    warnings.append(
                        f"slide {idx}: text overflows its box by "
                        f"{(cy - (y + h)) / DPI:.2f}\" — {runs[0].text[:42]!r}")

            if x < 0 or y < 0 or x + w > W + 2 or y + h > H + 2:
                warnings.append(f"slide {idx}: shape outside the slide "
                                f"({x},{y},{x+w},{y+h}) vs {W}x{H}")

        p = os.path.join(outdir, f"slide-{idx}.jpg")
        img.save(p, quality=88)
        written.append(p)
        print(p)

    print("\n" + ("no layout warnings" if not warnings else "WARNINGS:"))
    for w_ in warnings:
        print("  ! " + w_)
    return written


if __name__ == "__main__":
    render(sys.argv[1], sys.argv[2] if len(sys.argv) > 2 else "preview")
